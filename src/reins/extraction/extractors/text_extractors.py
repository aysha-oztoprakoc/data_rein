from typing import Any, Dict, Optional
from ..registry import registry
import os
import shutil
import sys
import json
import csv
import subprocess
from .base import BaseExtractor
import xml.etree.ElementTree as ET
from xml.dom import minidom
import re

try:
    import docx  # type: ignore
except ImportError:
    docx = None  # type: ignore
try:
    from bs4 import BeautifulSoup  # type: ignore
except ImportError:
    BeautifulSoup = None  # type: ignore
try:
    import ebooklib  # type: ignore
    from ebooklib import epub  # type: ignore
except ImportError:
    ebooklib = None  # type: ignore
    epub = None  # type: ignore
try:
    from striprtf.striprtf import rtf_to_text  # type: ignore
except ImportError:
    rtf_to_text = None  # type: ignore
try:
    import fitz  # type: ignore  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore
try:
    import openpyxl  # type: ignore
except ImportError:
    openpyxl = None  # type: ignore
try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None  # type: ignore

def remove_illegal_xml_chars(val: str) -> str:
    _illegal_xml_chars_RE = re.compile('[\x00-\x08\x0b\x0c\x0e-\x1F\uD800-\uDFFF\uFFFE\uFFFF]')
    ansi_escape = re.compile(r'\x1B(?:[@-Z\\-_]|\[[0-?]*[ -/]*[@-~])')
    val = ansi_escape.sub('', val)
    return _illegal_xml_chars_RE.sub('', val)

def save_as_xml(text_content: str, filepath: str, output_dir: str) -> str:
    out_path = os.path.join(output_dir, f"{os.path.basename(filepath)}.extracted.xml")
    content = remove_illegal_xml_chars(text_content)
    root = ET.Element("knowledge_document")
    meta = ET.SubElement(root, "metadata")
    title = ET.SubElement(meta, "title")
    title.text = remove_illegal_xml_chars(os.path.basename(filepath))
    path = ET.SubElement(meta, "path")
    path.text = remove_illegal_xml_chars(filepath)
    body = ET.SubElement(root, "content")
    body.text = content
    xml_str = minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(xml_str)
    return out_path


class PlainTextExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".txt", ".md"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "plaintext"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class JSONExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".json"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
            text = json.dumps(data, indent=2)
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "json"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class CSVExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".csv"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                text = "\n".join([", ".join(row) for row in reader])
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "csv"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class MinerUPDFExtractor(BaseExtractor):
    """
    Default PDF extractor: mineru (layout/table/OCR-aware) first, degrading to
    PyMuPDF (in-process, lighter) and finally the original `pdftotext`
    subprocess if either is unavailable or fails. mineru's layout-detection
    models are noticeably heavier than the other two paths, so large PDFs are
    routed straight to PyMuPDF to avoid pathological memory use on a
    resource-constrained box - never crashes, just degrades one step at a
    time per the harness's graceful-degradation convention.
    """
    SUPPORTED_FORMATS = [".pdf"]
    NODE = "amdy"

    # Page-count above which mineru's heavier layout inference is skipped in
    # favor of the lighter PyMuPDF path.
    MINERU_MAX_PAGES = 50
    # The FIRST call downloads ~7 model files and takes ~2min even on a
    # single-page PDF (measured); once cached, a 1-page doc takes ~15s. Sized
    # generously to survive that one-time cost rather than needlessly falling
    # back to PyMuPDF on the very first document processed.
    MINERU_TIMEOUT_S = 300

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        text, used = self._try_mineru(filepath, output_dir)
        if text is None:
            text, used = self._try_pymupdf(filepath)
        if text is None:
            text, used = self._try_pdftotext(filepath, output_dir)
        if text is None:
            return {"status": "error", "error": "all PDF extraction paths failed (mineru/pymupdf/pdftotext)"}
        out_path = save_as_xml(text, filepath, output_dir)
        return {"status": "success", "output_path": out_path, "metadata": {"format": "pdf", "extractor": used}}

    def _page_count(self, filepath: str) -> int:
        if fitz is None:
            return 0
        try:
            with fitz.open(filepath) as doc:
                return doc.page_count
        except Exception:
            return 0

    def _mineru_bin(self) -> Optional[str]:
        """Locate the `mineru` executable. It's often installed into the same
        venv as this interpreter (`<venv>/bin/mineru`), which isn't
        necessarily on $PATH when invoked via the `reins` console-script
        (its shebang points straight at the venv python) - so check next to
        sys.executable before giving up."""
        found = shutil.which("mineru")
        if found:
            return found
        candidate = os.path.join(os.path.dirname(sys.executable), "mineru")
        return candidate if os.path.isfile(candidate) and os.access(candidate, os.X_OK) else None

    def _try_mineru(self, filepath: str, output_dir: str):
        mineru_bin = self._mineru_bin()
        if mineru_bin is None:
            return None, None
        if self._page_count(filepath) > self.MINERU_MAX_PAGES:
            return None, None
        mineru_out = os.path.join(output_dir, f".mineru_{os.path.basename(filepath)}")
        try:
            os.makedirs(mineru_out, exist_ok=True)
            res = subprocess.run(
                # Pin to the "pipeline" backend explicitly: the CLI default
                # (hybrid-engine) pulls in VLM models, too heavy for an 8GB
                # VRAM / RAM-constrained box and liable to blow past the
                # timeout on first run downloading them.
                [mineru_bin, "-p", filepath, "-o", mineru_out, "-b", "pipeline"],
                capture_output=True, timeout=self.MINERU_TIMEOUT_S,
            )
            if res.returncode != 0:
                return None, None
            md_files = sorted(
                (os.path.join(root, f) for root, _, files in os.walk(mineru_out)
                 for f in files if f.endswith(".md")),
            )
            if not md_files:
                return None, None
            text = "\n\n".join(open(p, "r", encoding="utf-8", errors="replace").read() for p in md_files)
            return (text or None), "mineru"
        except Exception:
            return None, None
        finally:
            shutil.rmtree(mineru_out, ignore_errors=True)

    def _try_pymupdf(self, filepath: str):
        if fitz is None:
            return None, None
        try:
            with fitz.open(filepath) as doc:
                text = "\n".join(page.get_text() for page in doc)
            return (text or None), "pymupdf"
        except Exception:
            return None, None

    def _try_pdftotext(self, filepath: str, output_dir: str):
        out_path = os.path.join(output_dir, f"{os.path.basename(filepath)}.extracted.txt")
        try:
            res = subprocess.run(["pdftotext", filepath, out_path], capture_output=True)
            if res.returncode != 0:
                return None, None
            with open(out_path, 'r', encoding='utf-8') as f:
                text = f.read()
            os.remove(out_path)
            return (text or None), "pdftotext"
        except Exception:
            return None, None


class DocxExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".docx"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if docx is None:
            return {"status": "error", "error": "python-docx not installed"}
        try:
            doc = docx.Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "docx"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class HTMLExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".html", ".htm"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if BeautifulSoup is None:
            return {"status": "error", "error": "beautifulsoup4 not installed"}
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "html"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class XMLExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".xml"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        try:
            tree = ET.parse(filepath)
            root = tree.getroot()
            text = "".join(root.itertext())
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "xml"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class EpubExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".epub"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if epub is None or ebooklib is None or BeautifulSoup is None:
            return {"status": "error", "error": "ebooklib or bs4 not installed"}
        try:
            book = epub.read_epub(filepath)
            text = ""
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_body_content(), 'html.parser')
                    text += soup.get_text(separator="\n", strip=True) + "\n"
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "epub"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class RTFExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".rtf"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if rtf_to_text is None:
            return {"status": "error", "error": "striprtf not installed"}
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf = f.read()
            text = rtf_to_text(rtf)  # type: ignore[no-untyped-call]
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "rtf"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class XLSXExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".xlsx", ".xls"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if openpyxl is None:
            return {"status": "error", "error": "openpyxl not installed"}
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            chunks = []
            for sheet in wb.worksheets:
                chunks.append(f"--- SHEET: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    chunks.append(", ".join("" if v is None else str(v) for v in row))
            text = "\n".join(chunks)
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "xlsx"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


class PPTXExtractor(BaseExtractor):
    SUPPORTED_FORMATS = [".pptx"]
    NODE = "amdy"

    def extract(self, filepath: str, output_dir: str) -> Dict[str, Any]:
        if Presentation is None:
            return {"status": "error", "error": "python-pptx not installed"}
        try:
            prs = Presentation(filepath)
            chunks = []
            for i, slide in enumerate(prs.slides, start=1):
                chunks.append(f"--- SLIDE {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            text = "\n".join(chunks)
            out_path = save_as_xml(text, filepath, output_dir)
            return {"status": "success", "output_path": out_path, "metadata": {"format": "pptx"}}
        except Exception as e:
            return {"status": "error", "error": str(e)}


# Auto-register
registry.register(PlainTextExtractor)
registry.register(JSONExtractor)
registry.register(CSVExtractor)
registry.register(MinerUPDFExtractor)
registry.register(DocxExtractor)
registry.register(HTMLExtractor)
registry.register(XMLExtractor)
registry.register(EpubExtractor)
registry.register(RTFExtractor)
registry.register(XLSXExtractor)
registry.register(PPTXExtractor)
